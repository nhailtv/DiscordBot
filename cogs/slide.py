import base64
import glob
import os
import random
import re
import string
import requests
from urllib.parse import urlparse
import google.generativeai as genai
import discord
from discord.ext import commands
from pptx import Presentation
from pexels_api import API
from dotenv import load_dotenv

load_dotenv()

# Get the Gemini API Key from environment variable
GEMINI_API_KEY = os.getenv("GOOGLE_AI_KEY")  # Fetch the key from environment

# Create the cog class for generating PPT
class GeneratePPT(commands.Cog, name="generate_ppt"):
    def __init__(self, bot) -> None:
        self.bot = bot
        self.unique_image_name = None
        self.pexels_api_key = os.getenv("PEXELS_API_KEY")

        # Initialize the Pexels API
        self.pexels_client = API(self.pexels_api_key)

    def refresh_unique_image_name(self):
        self.unique_image_name = ''.join(random.choice(string.ascii_uppercase + string.ascii_lowercase + string.digits) for _ in range(16))


    async def generate_ppt(self, topic: str, slide_length: str):
        # Initialize the PowerPoint presentation
        slides_folder = r"C:\Users\nhail\Desktop\Discord_bot\Python-Discord-Bot-Template\cogs\Slides"
        root = Presentation(r"C:\Users\nhail\Desktop\Discord_bot\Python-Discord-Bot-Template\cogs\theme0.pptx")
        self.root = root  # Store the presentation in the instance attribute
        
        # Remove all existing slides before adding new ones
        self.remove_all_slides(root)

        # Configure Gemini API
        genai.configure(api_key=GEMINI_API_KEY)
        
        message = f"""Create an outline for a slideshow presentation on the topic of {topic} which is {slide_length}
        slides long. Make sure it is {slide_length} long.

        You are allowed to use the following slide types:
        Title Slide - (Title, Subtitle)
        Content Slide - (Title, Content)
        Image Slide - (Title, Content, Image)
        Thanks Slide - (Title)

        Put this tag before the Title Slide: [L_TS]
        Put this tag before the Content Slide: [L_CS]
        Put this tag before the Image Slide: [L_IS]
        Put this tag before the Thanks Slide: [L_THS]
        
        Put this tag before the Title: [TITLE]
        Put this tag after the Title: [/TITLE]
        Put this tag before the Subitle: [SUBTITLE]
        Put this tag after the Subtitle: [/SUBTITLE]
        Put this tag before the Content: [CONTENT]
        Put this tag after the Content: [/CONTENT]
        Put this tag before the Image: [IMAGE]
        Put this tag after the Image: [/IMAGE]

        Put "[SLIDEBREAK]" after each slide"""
        
        model = genai.GenerativeModel(model_name="gemini-1.5-pro")
        chat_session = model.start_chat(history=[])
        response = chat_session.send_message(message)

        list_of_slides = response.text.split("[SLIDEBREAK]")

        # Process each slide and create corresponding PowerPoint slide
        for slide in list_of_slides:
            slide_type = self.search_for_slide_type(slide)
            if slide_type == "[L_TS]":
                self.create_title_slide(self.find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]"),
                                        self.find_text_in_between_tags(str(slide), "[SUBTITLE]", "[/SUBTITLE]"))
            elif slide_type == "[L_CS]":
                self.create_title_and_content_slide(
                    "".join(self.find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")),
                    "".join(self.find_text_in_between_tags(str(slide), "[CONTENT]", "[/CONTENT]")))
            elif slide_type == "[L_IS]":
                self.create_title_and_content_and_image_slide("".join(self.find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")),
                                                            "".join(self.find_text_in_between_tags(str(slide), "[CONTENT]", "[/CONTENT]")) ,
                                                            "".join(self.find_text_in_between_tags(str(slide), "[IMAGE]", "[/IMAGE]")))
            elif slide_type == "[L_THS]":
                self.create_section_header_slide("".join(self.find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")))

        # Get a clean version of the title for the file name
        title = self.find_title()
        title = "".join(title.split(":")).strip()  # Clean title
        title = re.sub(r'[<>:"/\\|?*]', '', title)  # Remove invalid characters for filenames

        if not title:  # Fallback if no title found
            title = "Generated_Presentation"

        # Define the file path for saving the PPT
        slide_path = os.path.join(slides_folder, f"{title}.pptx")
        root.save(slide_path)

        return slide_path

    def remove_all_slides(self, presentation):
        # This function will remove all slides from the presentation
        for i in range(len(presentation.slides)-1, -1, -1):
            rId = presentation.slides._sldIdLst[i].rId
            presentation.part.drop_rel(rId)
            del presentation.slides._sldIdLst[i]

    def find_title(self):
        return self.root.slides[0].shapes.title.text

    def find_text_in_between_tags(self, text, start_tag, end_tag):
        start_pos = text.find(start_tag)
        end_pos = text.find(end_tag)
        result = []
        while start_pos > -1 and end_pos > -1:
            text_between_tags = text[start_pos + len(start_tag):end_pos]
            result.append(text_between_tags)
            start_pos = text.find(start_tag, end_pos + len(end_tag))
            end_pos = text.find(end_tag, start_pos)
        res1 = "".join(result)
        res2 = re.sub(r"\[IMAGE\].*?\[/IMAGE\]", '', res1)
        if len(result) > 0:
            return res2
        else:
            return ""

    def search_for_slide_type(self, text):
        tags = ["[L_TS]", "[L_CS]", "[L_IS]", "[L_THS]"]
        found_text = next((s for s in tags if s in text), None)
        return found_text

    def create_title_slide(self, title, subtitle):
        layout = self.root.slide_layouts[0]
        slide = self.root.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def create_section_header_slide(self, title):
        layout = self.root.slide_layouts[2]
        slide = self.root.slides.add_slide(layout)
        slide.shapes.title.text = title

    def create_title_and_content_slide(self, title, content):
        layout = self.root.slide_layouts[1]
        slide = self.root.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content

    def create_title_and_content_and_image_slide(self, title, content, image_query):
        layout = self.root.slide_layouts[8]
        slide = self.root.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[2].text = content
        self.refresh_unique_image_name()

        # Fetch images from Pexels using the API
        self.pexels_client.search(image_query, page=1, results_per_page=1)
        photos = self.pexels_client.get_entries()

        if photos:
            image_url = photos[0].original
            self.download_image(image_url)

            # Add the downloaded image to the slide
            img_path = f"C:\\Users\\nhail\\Desktop\\Discord_bot\\Python-Discord-Bot-Template\\cogs\\Slides\\p_{self.unique_image_name}.jpg"
            slide.shapes.add_picture(img_path, slide.placeholders[1].left, slide.placeholders[1].top,
                                     slide.placeholders[1].width, slide.placeholders[1].height)

    def download_image(self, image_url):
        # Download the image from Pexels and save it locally
        img_data = requests.get(image_url).content
        with open(f"C:\\Users\\nhail\\Desktop\\Discord_bot\\Python-Discord-Bot-Template\\cogs\\Slides\\p_{self.unique_image_name}.jpg", 'wb') as img_file:
            img_file.write(img_data)

    @commands.hybrid_command(
        name="slide",
        description="Generates a PPT based on the topic and slide length using Gemini API",
    )
    async def slide_command(self, ctx: commands.Context, topic: str, slide_length: int):
        """
        Generates a PowerPoint presentation based on the provided topic and slide length using the Gemini API.
        
        :param ctx: The application command context.
        :param topic: The topic for the PPT.
        :param slide_length: Number of slides for the PPT.
        """
        # Send an initial "working" message
        working_message = await ctx.send("The bot is working on generating the slides. Please wait...")

        try:
            # Generate the PowerPoint
            result = await self.generate_ppt(topic, str(slide_length))
            
            # Update the message to indicate that the process is done
            await working_message.edit(content=f"Done! The presentation is ready. You can download it now")
            
            # Send the slide file to Discord
            await ctx.send(file=discord.File(result))
            
        except Exception as e:
            # If something goes wrong, send an error message
            await working_message.edit(content=f"An error occurred while generating the slides: {e}")

# The setup function to add the cog to the bot
async def setup(bot) -> None:
    await bot.add_cog(GeneratePPT(bot))